"""
=============================================================
MODULE 1 - Step 1.1: Universal Document Scanner (v3.0 — Tank Edition)
=============================================================
Extracts text from 60+ file formats.

🛡️ v3.0 CRITICAL FIX: EasyOCR REMOVED!
   - EasyOCR requires torch + CUDA + ~500MB model download
   - Streamlit Cloud has 1GB RAM limit — EasyOCR crashes
   - Now uses ONLY pytesseract (system package from packages.txt)
   - pytesseract is faster (instant) and works everywhere

Supported Formats (60+):
- PDF (Text + Scanned via pytesseract)
- Images (JPG, PNG, BMP, TIFF, WebP, HEIC, GIF)
- Documents (DOCX, DOC, TXT, MD, RTF, ODT)
- Spreadsheets (XLSX, XLS, CSV, TSV)
- Presentations (PPTX)
- Data Formats (JSON, XML, HTML, HTM)
- eBooks (EPUB, MOBI)
- Code Files (PY, JS, TS, JAVA, C, CPP, GO, RS, etc.)

Security:
- F-01: Safe filename (blocks path traversal)
- F-07: File size enforcement (200MB)
- F-08: ZIP bomb protection (50MB per extracted member)
- F-14: Magic byte verification (anti-polyglot)
- F-15: Secure temp files (overwrite + delete)
- F-S5: All temp files use TEMP_DIR (no hardcoded paths)
- F-S6: OCR uses only pytesseract (Streamlit Cloud compatible)
=============================================================
"""

import os
import io
import re
import ssl
import json
import logging
import csv
import uuid
import zipfile
import tempfile
import shutil
import atexit
from pathlib import Path
from typing import Optional, List

from PIL import Image

logger = logging.getLogger(__name__)


# ============================================================
# 🛡️ FIX F-09: SSL verification — uses certifi's CA bundle
# ============================================================
def _get_safe_ssl_context():
    """Return a proper SSL context (verifies certificates)."""
    ctx = ssl.create_default_context()
    try:
        import certifi
        ctx.load_verify_locations(certifi.where())
    except ImportError:
        pass  # use system CA bundle
    return ctx


# ============================================================
# 🛡️ FIX F-15: Secure temp directory (auto-cleaned at exit)
# ============================================================
TEMP_DIR = Path(tempfile.mkdtemp(prefix='rag_app_'))
atexit.register(lambda: shutil.rmtree(TEMP_DIR, ignore_errors=True))


def _secure_cleanup(filepath: str) -> None:
    """Securely delete a file: overwrite with zeros, then remove."""
    try:
        if os.path.exists(filepath):
            size = os.path.getsize(filepath)
            with open(filepath, 'wb') as f:
                f.write(b'\x00' * min(size, 10 * 1024 * 1024))  # cap at 10MB
            os.remove(filepath)
    except OSError:
        pass  # best-effort cleanup


# ============================================================
# 🛡️ FIX F-01: Safe filename generation (blocks path traversal)
# ============================================================
def _safe_temp_path(original_name: str) -> Path:
    """Generate a safe temp path — no user-controlled traversal possible."""
    safe_name = Path(original_name).name  # 'a/b/c.pdf' -> 'c.pdf'
    safe_name = re.sub(r'[^A-Za-z0-9._-]', '_', safe_name)
    safe_name = safe_name.lstrip('.')  # strip leading dots
    safe_name = safe_name[:100] or 'upload'  # cap length, default if empty
    unique = f"{uuid.uuid4().hex[:8]}_{safe_name}"
    target = (TEMP_DIR / unique).resolve()
    # Defense in depth: verify target is inside TEMP_DIR
    if not str(target).startswith(str(TEMP_DIR.resolve())):
        raise ValueError("Path traversal detected")
    return target


# ============================================================
# 🛡️ FIX F-08: ZIP bomb protection (cap decompressed size)
# ============================================================
MAX_DECOMPRESSED_SIZE = 50 * 1024 * 1024  # 50MB per extracted file


def _safe_zip_extract(zf: zipfile.ZipFile, member: str) -> bytes:
    """Extract a single ZIP member with size cap (anti ZIP-bomb)."""
    info = zf.getinfo(member)
    if info.file_size > MAX_DECOMPRESSED_SIZE:
        raise ValueError(
            f"Decompressed size {info.file_size} exceeds limit {MAX_DECOMPRESSED_SIZE}"
        )
    data = zf.read(member)
    if len(data) > MAX_DECOMPRESSED_SIZE:
        raise ValueError("Actual decompressed size exceeds limit")
    return data


# ============================================================
# 🛡️ FIX F-14: Magic-byte verification (anti polyglot attack)
# ============================================================
MAGIC_BYTE_MAP = {
    'application/pdf': {'.pdf'},
    'image/jpeg': {'.jpg', '.jpeg'},
    'image/png': {'.png'},
    'image/bmp': {'.bmp'},
    'image/gif': {'.gif'},
    'image/webp': {'.webp'},
    'image/tiff': {'.tiff', '.tif'},
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': {'.docx'},
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': {'.xlsx'},
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': {'.pptx'},
    'application/epub+zip': {'.epub'},
    'application/vnd.oasis.opendocument.text': {'.odt'},
    'application/vnd.ms-excel': {'.xls'},
    'application/msword': {'.doc'},
    'application/vnd.ms-powerpoint': {'.ppt'},
}


def _verify_magic_bytes(file_path: str, claimed_ext: str) -> bool:
    """Verify file content matches claimed extension.

    Lenient mode for documents (PDFs sometimes mis-detected).
    Strict mode for images (PIL vulnerabilities are real).
    """
    try:
        import magic
        mime = magic.from_file(file_path, mime=True)
        allowed_exts = MAGIC_BYTE_MAP.get(mime, set())

        # Always accept these (too many false positives):
        if mime.startswith('text/') or mime == 'application/octet-stream':
            return True

        # For known document types, accept if extension matches
        if claimed_ext.lower() in allowed_exts:
            return True

        # Lenient: if extension is in our supported list AND magic detected
        # something plausible (not a totally different type), accept with warning
        if claimed_ext.lower() in {'.pdf', '.docx', '.xlsx', '.pptx', '.epub', '.odt'}:
            logger.warning(
                f"Magic byte mismatch for {claimed_ext} (detected: {mime}) — "
                f"accepting anyway (lenient mode for documents)"
            )
            return True

        # For images, be stricter (PIL vulnerabilities are real)
        logger.warning(f"Magic byte mismatch: claimed={claimed_ext}, detected={mime}")
        return False
    except ImportError:
        logger.debug("python-magic not installed — magic byte check skipped")
        return True  # fail-open if magic lib missing
    except Exception as e:
        logger.warning(f"Magic byte check failed ({e}) — accepting file")
        return True  # fail-open on any error


# ============================================================
# 📦 PDF AUTO-COMPRESSION (large files pe memory + speed optimization)
# ============================================================
PDF_COMPRESS_THRESHOLD_MB = 20  # Files above this size auto-compress


def _compress_pdf_if_large(file_path: str, threshold_mb: int = PDF_COMPRESS_THRESHOLD_MB) -> str:
    """Compress large PDF to reduce memory + speed up processing.

    Returns:
        Path to compressed PDF (or original path if compression skipped/failed).
    """
    try:
        file_size = os.path.getsize(file_path)
        if file_size < threshold_mb * 1024 * 1024:
            return file_path  # No need to compress

        try:
            import fitz  # PyMuPDF
        except ImportError:
            logger.warning(
                "⚠️ PyMuPDF not installed — cannot auto-compress large PDF. "
                "Install with: pip install PyMuPDF"
            )
            return file_path

        logger.info(f"📦 Auto-compressing large PDF ({file_size/1024/1024:.1f}MB)...")

        doc = fitz.open(file_path)
        if doc.is_encrypted:
            doc.close()
            logger.warning("⚠️ Cannot compress encrypted PDF, using original")
            return file_path

        compressed_path = str(TEMP_DIR / f"compressed_{uuid.uuid4().hex[:8]}_{Path(file_path).name}")
        doc.save(
            compressed_path,
            garbage=4,         # Max garbage collection
            deflate=True,      # Compress all streams
            clean=True,        # Sanitize content
            linear=True,       # Linearize for faster access
        )
        doc.close()

        compressed_size = os.path.getsize(compressed_path)
        ratio = (1 - compressed_size / file_size) * 100

        # If compression made file LARGER (rare), use original
        if compressed_size >= file_size * 0.95:
            _secure_cleanup(compressed_path)
            logger.info(
                f"ℹ️ Compression didn't help (original={file_size/1024/1024:.1f}MB, "
                f"compressed={compressed_size/1024/1024:.1f}MB) — using original"
            )
            return file_path

        logger.info(
            f"✅ PDF compressed: {file_size/1024/1024:.1f}MB → "
            f"{compressed_size/1024/1024:.1f}MB ({ratio:.0f}% smaller)"
        )

        # Stage 2: If still very large, downsample images
        if compressed_size > 30 * 1024 * 1024:
            try:
                further_compressed = _downsample_pdf_images(compressed_path)
                if further_compressed != compressed_path:
                    _secure_cleanup(compressed_path)
                    return further_compressed
            except Exception as e:
                logger.warning(f"Image downsampling skipped: {e}")

        return compressed_path

    except Exception as e:
        logger.warning(f"⚠️ PDF compression failed ({e}) — using original")
        return file_path


def _downsample_pdf_images(file_path: str, target_dpi: int = 150) -> str:
    """Downsample high-res images inside PDF to reduce size further."""
    import fitz

    doc = fitz.open(file_path)
    if doc.is_encrypted:
        doc.close()
        return file_path

    max_dimension = 2000
    images_downsampled = 0

    for page in doc:
        image_list = page.get_images(full=True)
        for img_info in image_list:
            xref = img_info[0]
            try:
                pix = fitz.Pixmap(doc, xref)
                if pix.width <= max_dimension and pix.height <= max_dimension:
                    pix = None
                    continue
                # Calculate new dimensions
                if pix.width >= pix.height:
                    new_w = max_dimension
                    new_h = int(pix.height * (max_dimension / pix.width))
                else:
                    new_h = max_dimension
                    new_w = int(pix.width * (max_dimension / pix.height))
                new_pix = fitz.Pixmap(pix, new_w, new_h)
                new_pix = None
                pix = None
                images_downsampled += 1
            except Exception:
                continue

    if images_downsampled == 0:
        doc.close()
        return file_path

    output_path = str(TEMP_DIR / f"downsampled_{uuid.uuid4().hex[:8]}_{Path(file_path).name}")
    doc.save(output_path, garbage=4, deflate=True, clean=True)
    doc.close()

    new_size = os.path.getsize(output_path)
    old_size = os.path.getsize(file_path)
    logger.info(
        f"🖼️ Downsampled {images_downsampled} images: "
        f"{old_size/1024/1024:.1f}MB → {new_size/1024/1024:.1f}MB"
    )

    if new_size >= old_size:
        _secure_cleanup(output_path)
        return file_path

    return output_path


# ============================================================
# 🛡️ v3.0 FIX F-S6: Tesseract-only OCR (NO EasyOCR)
# ============================================================
# Previously: EasyOCR was used as primary, pytesseract as fallback.
# Problem: EasyOCR crashes on Streamlit Cloud (1GB RAM, no CUDA, slow downloads).
# Now: pytesseract is the ONLY OCR engine. It's:
#   - Instant (no model download)
#   - Lightweight (uses system tesseract-ocr from packages.txt)
#   - Works everywhere (Linux, Mac, Windows, Streamlit Cloud, Docker)
#   - Supports 100+ languages via apt install tesseract-ocr-LANG

def _check_tesseract_available() -> bool:
    """Check if pytesseract + tesseract binary are available."""
    try:
        import pytesseract
        version = pytesseract.get_tesseract_version()
        logger.info(f"✅ Tesseract OCR available (v{version})")
        return True
    except ImportError:
        logger.error(
            "❌ pytesseract not installed! "
            "Install with: pip install pytesseract"
        )
        return False
    except Exception as e:
        logger.error(
            f"❌ Tesseract binary not found! "
            f"Install system package: apt install tesseract-ocr  "
            f"(On Streamlit Cloud: add 'tesseract-ocr' to packages.txt). "
            f"Error: {e}"
        )
        return False


def _get_tesseract_languages() -> List[str]:
    """Get list of available Tesseract languages on this system."""
    try:
        import pytesseract
        langs = pytesseract.get_languages()
        return langs or ['eng']
    except Exception:
        return ['eng']


def _tesseract_ocr_image(img: Image.Image, languages: List[str]) -> str:
    """Run Tesseract OCR on a PIL Image.

    Args:
        img: PIL Image to OCR
        languages: List of language codes (e.g., ['eng'], ['eng', 'hin'])

    Returns:
        Extracted text (empty string if OCR fails or no text found)
    """
    try:
        import pytesseract
        # Filter to only available languages
        available = set(_get_tesseract_languages())
        # Convert our codes to tesseract codes
        tesseract_langs = []
        for lang in languages:
            # 'en' -> 'eng', 'hi' -> 'hin', etc.
            tesseract_code = {'en': 'eng', 'hi': 'hin', 'ur': 'urd',
                             'ar': 'ara', 'zh': 'chi_sim', 'ja': 'jpn'}.get(lang, lang)
            if tesseract_code in available:
                tesseract_langs.append(tesseract_code)

        if not tesseract_langs:
            tesseract_langs = ['eng']  # Fallback to English

        lang_str = '+'.join(tesseract_langs)
        text = pytesseract.image_to_string(img, lang=lang_str)
        return text if text and len(text.strip()) > 5 else ""
    except Exception as e:
        logger.debug(f"Tesseract OCR skipped: {e}")
        return ""


# ============================================================
# Supported Formats Registry — 60+ Formats!
# ============================================================

SUPPORTED_FORMATS = {
    # PDFs
    '.pdf':   {'name': 'PDF Document',        'method': 'pdf',     'icon': '📄'},
    # Images
    '.jpg':   {'name': 'JPEG Image',           'method': 'image',   'icon': '🖼️'},
    '.jpeg':  {'name': 'JPEG Image',           'method': 'image',   'icon': '🖼️'},
    '.png':   {'name': 'PNG Image',            'method': 'image',   'icon': '🖼️'},
    '.bmp':   {'name': 'BMP Image',            'method': 'image',   'icon': '🖼️'},
    '.tiff':  {'name': 'TIFF Image',           'method': 'image',   'icon': '🖼️'},
    '.tif':   {'name': 'TIFF Image',           'method': 'image',   'icon': '🖼️'},
    '.webp':  {'name': 'WebP Image',           'method': 'image',   'icon': '🖼️'},
    '.gif':   {'name': 'GIF Image',            'method': 'image',   'icon': '🖼️'},
    '.heic':  {'name': 'HEIC Image (Apple)',   'method': 'image',   'icon': '🖼️'},
    # Word Documents
    '.docx':  {'name': 'Word Document',        'method': 'docx',   'icon': '📝'},
    '.doc':   {'name': 'Word Document (Old)',   'method': 'docx',   'icon': '📝'},
    '.rtf':   {'name': 'Rich Text Format',     'method': 'rtf',    'icon': '📝'},
    '.odt':   {'name': 'OpenDocument Text',    'method': 'odt',    'icon': '📝'},
    # Text Files
    '.txt':   {'name': 'Text File',            'method': 'text',   'icon': '📃'},
    '.md':    {'name': 'Markdown File',        'method': 'text',   'icon': '📋'},
    '.log':   {'name': 'Log File',             'method': 'text',   'icon': '📃'},
    '.ini':   {'name': 'Config File (INI)',    'method': 'text',   'icon': '⚙️'},
    '.yaml':  {'name': 'YAML File',            'method': 'text',   'icon': '📋'},
    '.yml':   {'name': 'YAML File',            'method': 'text',   'icon': '📋'},
    '.toml':  {'name': 'TOML Config',          'method': 'text',   'icon': '⚙️'},
    # Code / Programming Files
    '.py':    {'name': 'Python Script',         'method': 'text',   'icon': '🐍'},
    '.js':    {'name': 'JavaScript File',       'method': 'text',   'icon': '⚡'},
    '.ts':    {'name': 'TypeScript File',       'method': 'text',   'icon': '🔷'},
    '.java':  {'name': 'Java Source',           'method': 'text',   'icon': '☕'},
    '.c':     {'name': 'C Source',              'method': 'text',   'icon': '⚙️'},
    '.cpp':   {'name': 'C++ Source',            'method': 'text',   'icon': '⚙️'},
    '.h':     {'name': 'C/C++ Header',          'method': 'text',   'icon': '⚙️'},
    '.cs':    {'name': 'C# Source',             'method': 'text',   'icon': '⚙️'},
    '.go':    {'name': 'Go Source',              'method': 'text',   'icon': '🔵'},
    '.rs':    {'name': 'Rust Source',            'method': 'text',   'icon': '🦀'},
    '.rb':    {'name': 'Ruby Script',            'method': 'text',   'icon': '💎'},
    '.php':   {'name': 'PHP Script',             'method': 'text',   'icon': '🐘'},
    '.sh':    {'name': 'Shell Script',            'method': 'text',   'icon': '🖥️'},
    '.bat':   {'name': 'Batch Script',            'method': 'text',   'icon': '🖥️'},
    '.ps1':   {'name': 'PowerShell Script',       'method': 'text',   'icon': '🖥️'},
    '.sql':   {'name': 'SQL Script',              'method': 'text',   'icon': '🗄️'},
    '.css':   {'name': 'CSS Stylesheet',          'method': 'text',   'icon': '🎨'},
    '.scss':  {'name': 'SCSS Stylesheet',         'method': 'text',   'icon': '🎨'},
    '.r':     {'name': 'R Script',                'method': 'text',   'icon': '📊'},
    '.swift': {'name': 'Swift Source',            'method': 'text',   'icon': '🍎'},
    '.kt':    {'name': 'Kotlin Source',           'method': 'text',   'icon': '🟣'},
    '.dart':  {'name': 'Dart Source',             'method': 'text',   'icon': '🎯'},
    '.lua':   {'name': 'Lua Script',              'method': 'text',   'icon': '🌙'},
    '.scala': {'name': 'Scala Source',            'method': 'text',   'icon': '🔴'},
    '.vue':   {'name': 'Vue Component',           'method': 'text',   'icon': '💚'},
    '.jsx':   {'name': 'React JSX',              'method': 'text',   'icon': '⚛️'},
    '.tsx':   {'name': 'React TSX',              'method': 'text',   'icon': '⚛️'},
    # Spreadsheets
    '.xlsx':  {'name': 'Excel Spreadsheet',    'method': 'xlsx',   'icon': '📈'},
    '.xls':   {'name': 'Excel (Old)',          'method': 'xlsx',   'icon': '📈'},
    '.csv':   {'name': 'CSV Data',             'method': 'csv',    'icon': '📊'},
    '.tsv':   {'name': 'TSV Data',             'method': 'csv',    'icon': '📊'},
    # Presentations
    '.pptx':  {'name': 'PowerPoint',           'method': 'pptx',   'icon': '📽️'},
    # Data Formats
    '.json':  {'name': 'JSON Data',            'method': 'json',   'icon': '🔗'},
    '.xml':   {'name': 'XML Document',         'method': 'xml',    'icon': '📰'},
    '.html':  {'name': 'HTML Page',            'method': 'html',   'icon': '🌐'},
    '.htm':   {'name': 'HTML Page (Old)',      'method': 'html',   'icon': '🌐'},
    # eBooks
    '.epub':  {'name': 'eBook (EPUB)',         'method': 'epub',   'icon': '📖'},
    '.mobi':  {'name': 'eBook (MOBI)',         'method': 'mobi',   'icon': '📖'},
}

# Extensions list for Streamlit file_uploader
ACCEPTED_EXTENSIONS = sorted(list(SUPPORTED_FORMATS.keys()))


# ============================================================
# OCR Scanner Class
# ============================================================
class OCRScanner:
    """
    Universal Document Scanner — 60+ Formats!

    🛡️ v3.0: Uses ONLY Tesseract OCR (no EasyOCR).
    - Fast: instant, no model download
    - Lightweight: uses system tesseract-ocr binary
    - Cloud-compatible: works on Streamlit Cloud, Docker, etc.
    """

    SUPPORTED_LANGUAGES = {
        'en': 'English', 'hi': 'Hindi', 'ur': 'Urdu',
        'ar': 'Arabic', 'zh': 'Chinese', 'ja': 'Japanese',
    }

    def __init__(self, languages: list = None):
        self.languages = languages or ['en']
        self._tesseract_available = None
        self._available_langs = None

    @property
    def tesseract_available(self) -> bool:
        """Check if pytesseract + tesseract binary are available."""
        if self._tesseract_available is None:
            self._tesseract_available = _check_tesseract_available()
        return self._tesseract_available

    @property
    def available_langs(self) -> List[str]:
        """Get list of available Tesseract languages."""
        if self._available_langs is None:
            self._available_langs = _get_tesseract_languages()
        return self._available_langs

    @property
    def ocr_available(self) -> bool:
        """Check if OCR is available (just tesseract now)."""
        return self.tesseract_available

    # ============================================================
    # MAIN ENTRY POINT
    # ============================================================

    def extract_from_uploaded_file(self, uploaded_file) -> str:
        """Streamlit uploaded file se text extract karo. Supports 60+ formats!

        🛡️ Hardened with:
        - F-01: Safe filename (path traversal blocked)
        - F-07: File size enforcement (200MB)
        - F-14: Magic-byte verification (polyglot blocked)
        - F-15: Secure temp file handling (zero-then-delete)
        - F-S5: All temp files use TEMP_DIR
        """
        # 🛡️ F-07: File size enforcement
        try:
            file_size = uploaded_file.size
        except (AttributeError, TypeError):
            file_size = len(uploaded_file.getvalue()) if hasattr(uploaded_file, 'getvalue') else 0

        MAX_SIZE = 200 * 1024 * 1024  # 200MB
        if file_size > MAX_SIZE:
            logger.error(f"❌ File too large: {file_size/1024/1024:.1f}MB > {MAX_SIZE/1024/1024:.0f}MB")
            return ""
        if file_size == 0:
            logger.error("❌ Empty file (0 bytes)")
            return ""

        # 🛡️ F-01: Safe temp path
        temp_path = _safe_temp_path(uploaded_file.name)

        try:
            with open(temp_path, "wb") as f:
                f.write(uploaded_file.getbuffer() if hasattr(uploaded_file, 'getbuffer') else uploaded_file.read())

            file_ext = os.path.splitext(uploaded_file.name)[1].lower()
            fmt_info = SUPPORTED_FORMATS.get(file_ext)

            if not fmt_info:
                logger.error(f"❌ Unsupported file type: {file_ext}")
                return ""

            # 🛡️ F-14: Magic byte verification
            if not _verify_magic_bytes(str(temp_path), file_ext):
                logger.error(f"❌ Magic byte mismatch for {uploaded_file.name} (claimed {file_ext})")
                return ""

            method = fmt_info['method']
            logger.info(f"{fmt_info['icon']} Processing {fmt_info['name']}: {uploaded_file.name}")

            # Route to correct extraction method
            extractors = {
                'pdf': self.extract_from_pdf,
                'image': self.extract_from_image,
                'docx': self.extract_from_docx,
                'rtf': self.extract_from_rtf,
                'odt': self.extract_from_odt,
                'text': self.extract_from_text,
                'xlsx': self.extract_from_xlsx,
                'csv': self.extract_from_csv,
                'pptx': self.extract_from_pptx,
                'json': self.extract_from_json,
                'xml': self.extract_from_xml,
                'html': self.extract_from_html,
                'epub': self.extract_from_epub,
                'mobi': self.extract_from_mobi,
            }

            extractor = extractors.get(method)
            if extractor:
                return extractor(str(temp_path))
            return ""

        except (OSError, IOError) as e:
            logger.error(f"❌ Filesystem error: {e}", exc_info=True)
            return ""
        except Exception as e:
            logger.exception(f"❌ File processing failed: {e}")
            return ""
        finally:
            # 🛡️ F-15: Secure cleanup (overwrite + delete)
            _secure_cleanup(str(temp_path))

    # ============================================================
    # PDF EXTRACTION
    # ============================================================

    def extract_from_pdf(self, file_path: str) -> str:
        """PDF se text nikalo — PyPDF2 pehle, Tesseract OCR fallback.

        📦 Large PDFs (>20MB) auto-compress hote hain pehle (memory + speed).
        """
        all_text = ""

        # 📦 STEP 0: Auto-compress large PDF (transparent to caller)
        original_size = os.path.getsize(file_path)
        actual_path = _compress_pdf_if_large(file_path)
        if actual_path != file_path:
            compressed_size = os.path.getsize(actual_path)
            logger.info(
                f"📦 Using compressed PDF: "
                f"{original_size/1024/1024:.1f}MB → {compressed_size/1024/1024:.1f}MB"
            )
        else:
            compressed_size = original_size

        try:
            # STEP 1: Try PyPDF2 for text extraction
            from PyPDF2 import PdfReader
            pdf_reader = PdfReader(actual_path)
            total_pages = len(pdf_reader.pages)
            logger.info(f"📄 PDF loaded: {total_pages} pages")

            for i, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                except Exception as e:
                    logger.warning(f"⚠️ Page {i+1} text extraction failed: {e}")
                    page_text = ""

                if page_text and len(page_text.strip()) > 30:
                    all_text += f"\n--- Page {i+1} ---\n{page_text}\n"
                else:
                    # STEP 2: OCR fallback for image-only pages
                    ocr_text = self._ocr_pdf_page(actual_path, i)
                    if ocr_text:
                        all_text += f"\n--- Page {i+1} (OCR) ---\n{ocr_text}\n"

        except ImportError:
            logger.warning("⚠️ PyPDF2 not installed, trying pure OCR")
            all_text = self._ocr_entire_pdf(actual_path)
        except Exception as e:
            logger.error(f"❌ PDF extraction failed: {e}", exc_info=True)
            all_text = self._ocr_entire_pdf(actual_path)
        finally:
            # 📦 Clean up compressed file if it was created
            if actual_path != file_path:
                _secure_cleanup(actual_path)

        return self._clean_text(all_text)

    # ============================================================
    # IMAGE EXTRACTION — Tesseract Only (v3.0)
    # ============================================================

    def extract_from_image(self, file_path: str) -> str:
        """Image se text nikalo using Tesseract OCR.

        🛡️ v3.0: EasyOCR removed. Tesseract is fast + cloud-compatible.
        """
        if not self.tesseract_available:
            logger.error(
                "❌ Tesseract OCR not available! "
                "Install: apt install tesseract-ocr (or add to packages.txt on Streamlit Cloud)"
            )
            return ""

        try:
            img = Image.open(file_path)

            # Handle multi-page TIFFs
            if hasattr(img, 'n_frames') and img.n_frames > 1:
                all_text = ""
                for frame_num in range(img.n_frames):
                    img.seek(frame_num)
                    # Convert to RGB if necessary (Tesseract prefers RGB)
                    if img.mode not in ('RGB', 'L'):
                        img_frame = img.convert('RGB')
                    else:
                        img_frame = img
                    text = _tesseract_ocr_image(img_frame, self.languages)
                    if text:
                        all_text += f"\n--- Frame {frame_num+1} ---\n{text}\n"
                return self._clean_text(all_text)

            # Single image
            if img.mode not in ('RGB', 'L'):
                img = img.convert('RGB')
            text = _tesseract_ocr_image(img, self.languages)
            if text:
                logger.info(f"⚡ Tesseract OCR: {len(text)} chars")
            else:
                logger.warning("⚠️ Tesseract returned no text (image may have no text)")
            return self._clean_text(text)

        except Exception as e:
            logger.error(f"❌ Image OCR failed: {e}")
            return ""

    # ============================================================
    # DOCX EXTRACTION
    # ============================================================

    def extract_from_docx(self, file_path: str) -> str:
        """Word Document (.docx) se text nikalo."""
        try:
            from docx import Document
            doc = Document(file_path)

            all_text = ""
            for para in doc.paragraphs:
                if para.text.strip():
                    all_text += para.text + "\n"

            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        all_text += row_text + "\n"

            logger.info(f"📝 DOCX extracted: {len(all_text)} chars")
            return self._clean_text(all_text)

        except ImportError:
            logger.error("❌ python-docx not installed! Run: pip install python-docx")
            return ""
        except Exception as e:
            logger.error(f"❌ DOCX extraction failed: {e}")
            return ""

    # ============================================================
    # RTF EXTRACTION
    # ============================================================

    def extract_from_rtf(self, file_path: str) -> str:
        """Rich Text Format (.rtf) se text nikalo."""
        try:
            try:
                from striprtf.striprtf import rtf_to_text
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    rtf_content = f.read()
                text = rtf_to_text(rtf_content)
                if text and len(text.strip()) > 5:
                    logger.info(f"📝 RTF extracted: {len(text)} chars")
                    return self._clean_text(text)
            except ImportError:
                pass

            # Fallback: basic RTF stripping using regex
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            text = re.sub(r'\\[a-z]+\d*\s?', '', content)
            text = re.sub(r'[{}]', '', text)
            text = re.sub(r'\\[^a-z]', '', text)
            logger.info(f"📝 RTF extracted (basic): {len(text)} chars")
            return self._clean_text(text)

        except Exception as e:
            logger.error(f"❌ RTF extraction failed: {e}")
            return ""

    # ============================================================
    # ODT EXTRACTION (OpenDocument Text)
    # ============================================================

    def extract_from_odt(self, file_path: str) -> str:
        """OpenDocument Text (.odt) se text nikalo — LibreOffice format.

        🛡️ F-08: Uses _safe_zip_extract (ZIP bomb protection).
        """
        try:
            import zipfile
            with zipfile.ZipFile(file_path, 'r') as zf:
                if 'content.xml' in zf.namelist():
                    xml_content = _safe_zip_extract(zf, 'content.xml').decode('utf-8', errors='ignore')
                    text = re.sub(r'<[^>]+>', ' ', xml_content)
                    text = re.sub(r'\s+', ' ', text)
                    logger.info(f"📝 ODT extracted: {len(text)} chars")
                    return self._clean_text(text)
            return ""
        except ValueError as e:
            logger.error(f"❌ ODT ZIP bomb detected: {e}")
            return ""
        except Exception as e:
            logger.error(f"❌ ODT extraction failed: {e}", exc_info=True)
            return ""

    # ============================================================
    # TEXT FILES (TXT, MD, LOG, INI, YAML, TOML, code)
    # ============================================================

    def extract_from_text(self, file_path: str) -> str:
        """TXT, MD, LOG, INI, YAML, YML, TOML, code files se direct text padho."""
        try:
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
            except UnicodeDecodeError:
                with open(file_path, 'r', encoding='latin-1') as f:
                    text = f.read()

            logger.info(f"📃 Text file read: {len(text)} chars")
            return self._clean_text(text)
        except Exception as e:
            logger.error(f"❌ Text file read failed: {e}")
            return ""

    # ============================================================
    # CSV / TSV EXTRACTION
    # ============================================================

    def extract_from_csv(self, file_path: str) -> str:
        """CSV aur TSV files se structured text nikalo."""
        try:
            file_ext = os.path.splitext(file_path)[1].lower()
            delimiter = '\t' if file_ext == '.tsv' else ','
            all_text = ""

            for encoding in ['utf-8', 'latin-1', 'cp1252']:
                try:
                    with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                        reader = csv.reader(f, delimiter=delimiter)
                        for i, row in enumerate(reader):
                            if i == 0:
                                all_text += " | ".join(row) + "\n"
                            else:
                                row_text = " | ".join(cell.strip() for cell in row if cell.strip())
                                if row_text.strip(" |"):
                                    all_text += row_text + "\n"
                    break
                except Exception:
                    continue

            logger.info(f"📊 CSV/TSV extracted: {len(all_text)} chars")
            return self._clean_text(all_text)

        except Exception as e:
            logger.error(f"❌ CSV/TSV extraction failed: {e}")
            return ""

    # ============================================================
    # XLSX EXTRACTION
    # ============================================================

    def extract_from_xlsx(self, file_path: str) -> str:
        """Excel Spreadsheet se text content nikalo."""
        try:
            from openpyxl import load_workbook
            wb = load_workbook(file_path, read_only=True, data_only=True)

            all_text = ""
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                all_text += f"\n=== Sheet: {sheet_name} ===\n"

                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip(" |"):
                        all_text += row_text + "\n"

            wb.close()
            logger.info(f"📈 XLSX extracted: {len(all_text)} chars")
            return self._clean_text(all_text)

        except ImportError:
            logger.error("❌ openpyxl not installed! Run: pip install openpyxl")
            return ""
        except Exception as e:
            logger.error(f"❌ XLSX extraction failed: {e}")
            return ""

    # ============================================================
    # PPTX EXTRACTION
    # ============================================================

    def extract_from_pptx(self, file_path: str) -> str:
        """PowerPoint Presentation se text nikalo."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)

            all_text = ""
            for i, slide in enumerate(prs.slides):
                all_text += f"\n=== Slide {i+1} ===\n"

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        all_text += shape.text + "\n"

                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join(cell.text.strip() for cell in row.cells)
                            if row_text.strip():
                                all_text += row_text + "\n"

            logger.info(f"📽️ PPTX extracted: {len(all_text)} chars")
            return self._clean_text(all_text)

        except ImportError:
            logger.error("❌ python-pptx not installed! Run: pip install python-pptx")
            return ""
        except Exception as e:
            logger.error(f"❌ PPTX extraction failed: {e}")
            return ""

    # ============================================================
    # JSON EXTRACTION
    # ============================================================

    def extract_from_json(self, file_path: str) -> str:
        """JSON file se structured text nikalo."""
        try:
            data = None
            for encoding in ['utf-8', 'latin-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        data = json.load(f)
                    break
                except (UnicodeDecodeError, json.JSONDecodeError):
                    continue

            if data is None:
                return self.extract_from_text(file_path)

            text = self._json_to_text(data)
            logger.info(f"🔗 JSON extracted: {len(text)} chars")
            return self._clean_text(text)

        except Exception as e:
            logger.error(f"❌ JSON extraction failed: {e}")
            return self.extract_from_text(file_path)

    def _json_to_text(self, data, prefix="") -> str:
        """Recursively convert JSON to readable text."""
        parts = []
        if isinstance(data, dict):
            for key, value in data.items():
                if isinstance(value, (dict, list)):
                    parts.append(f"{prefix}{key}:")
                    parts.append(self._json_to_text(value, prefix + "  "))
                else:
                    parts.append(f"{prefix}{key}: {value}")
        elif isinstance(data, list):
            for i, item in enumerate(data):
                if isinstance(item, (dict, list)):
                    parts.append(f"{prefix}Item {i+1}:")
                    parts.append(self._json_to_text(item, prefix + "  "))
                else:
                    parts.append(f"{prefix}{item}")
        else:
            parts.append(f"{prefix}{data}")
        return "\n".join(parts)

    # ============================================================
    # XML EXTRACTION
    # ============================================================

    def extract_from_xml(self, file_path: str) -> str:
        """XML file se text content nikalo."""
        try:
            import xml.etree.ElementTree as ET
            tree = ET.parse(file_path)
            root = tree.getroot()

            text_parts = []
            self._xml_to_text(root, text_parts)
            all_text = "\n".join(text_parts)

            logger.info(f"📰 XML extracted: {len(all_text)} chars")
            return self._clean_text(all_text)

        except Exception as e:
            logger.error(f"❌ XML extraction failed: {e}")
            return self.extract_from_text(file_path)

    def _xml_to_text(self, element, parts: list, depth=0):
        """Recursively extract text from XML elements."""
        if len(element) > 0:
            parts.append(f"{'  ' * depth}[{element.tag}]")

        if element.text and element.text.strip():
            parts.append(f"{'  ' * depth}{element.text.strip()}")

        if element.tail and element.tail.strip():
            parts.append(f"{'  ' * depth}{element.tail.strip()}")

        for child in element:
            self._xml_to_text(child, parts, depth + 1)

    # ============================================================
    # HTML EXTRACTION
    # ============================================================

    def extract_from_html(self, file_path: str) -> str:
        """HTML file se visible text nikalo (tags remove karo)."""
        try:
            try:
                from bs4 import BeautifulSoup
                html_content = None
                for encoding in ['utf-8', 'latin-1']:
                    try:
                        with open(file_path, 'r', encoding=encoding) as f:
                            html_content = f.read()
                        break
                    except UnicodeDecodeError:
                        continue

                if html_content is None:
                    return ""

                soup = BeautifulSoup(html_content, 'html.parser')

                # Remove script and style elements
                for script in soup(["script", "style", "noscript"]):
                    script.decompose()

                text = soup.get_text(separator='\n')
                lines = [line.strip() for line in text.splitlines() if line.strip()]
                text = "\n".join(lines)

                logger.info(f"🌐 HTML extracted (BeautifulSoup): {len(text)} chars")
                return self._clean_text(text)
            except ImportError:
                pass

            # Fallback: regex-based HTML stripping
            html_content = ""
            for encoding in ['utf-8', 'latin-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        html_content = f.read()
                    break
                except UnicodeDecodeError:
                    continue

            text = re.sub(r'<script[^>]*>.*?</script>', '', html_content, flags=re.DOTALL | re.IGNORECASE)
            text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)
            text = re.sub(r'<[^>]+>', ' ', text)
            text = re.sub(r'&nbsp;', ' ', text)
            text = re.sub(r'&amp;', '&', text)
            text = re.sub(r'&lt;', '<', text)
            text = re.sub(r'&gt;', '>', text)
            text = re.sub(r'&quot;', '"', text)
            text = re.sub(r'&#\d+;', '', text)

            logger.info(f"🌐 HTML extracted (regex): {len(text)} chars")
            return self._clean_text(text)

        except Exception as e:
            logger.error(f"❌ HTML extraction failed: {e}")
            return ""

    # ============================================================
    # EPUB EXTRACTION (eBooks)
    # ============================================================

    def extract_from_epub(self, file_path: str) -> str:
        """EPUB eBook se text nikalo."""
        try:
            try:
                import ebooklib
                from ebooklib import epub as epublib
                from bs4 import BeautifulSoup

                book = epublib.read_epub(file_path)
                all_text = ""

                for i, item in enumerate(book.get_items_of_type(ebooklib.ITEM_DOCUMENT)):
                    soup = BeautifulSoup(item.get_content(), 'html.parser')

                    for script in soup(["script", "style"]):
                        script.decompose()

                    text = soup.get_text(separator='\n')
                    lines = [line.strip() for line in text.splitlines() if line.strip()]
                    page_text = "\n".join(lines)

                    if page_text:
                        all_text += f"\n--- Chapter {i+1} ---\n{page_text}\n"

                logger.info(f"📖 EPUB extracted: {len(all_text)} chars")
                return self._clean_text(all_text)
            except ImportError:
                pass

            # Fallback: EPUB is a ZIP file, extract HTML content
            import zipfile
            all_text = ""
            with zipfile.ZipFile(file_path, 'r') as zf:
                for name in sorted(zf.namelist()):
                    if name.endswith(('.html', '.xhtml', '.htm')):
                        try:
                            html_content = _safe_zip_extract(zf, name).decode('utf-8', errors='ignore')
                        except ValueError:
                            logger.warning(f"Skipping oversized EPUB member: {name}")
                            continue
                        text = re.sub(r'<[^>]+>', ' ', html_content)
                        text = re.sub(r'\s+', ' ', text)
                        if len(text.strip()) > 20:
                            all_text += text.strip() + "\n"

            logger.info(f"📖 EPUB extracted (zip): {len(all_text)} chars")
            return self._clean_text(all_text)

        except Exception as e:
            logger.error(f"❌ EPUB extraction failed: {e}")
            return ""

    # ============================================================
    # MOBI EXTRACTION (Kindle eBooks)
    # ============================================================

    def extract_from_mobi(self, file_path: str) -> str:
        """MOBI eBook se text nikalo."""
        try:
            try:
                import mobi
                tempdir, filepath = mobi.extract(file_path)
                if os.path.exists(filepath):
                    result = self.extract_from_html(filepath)
                    import shutil
                    shutil.rmtree(tempdir, ignore_errors=True)
                    return result
            except ImportError:
                pass

            logger.warning("⚠️ MOBI support limited. Install 'mobi' package for better results: pip install mobi")
            return self.extract_from_text(file_path)

        except Exception as e:
            logger.error(f"❌ MOBI extraction failed: {e}")
            return ""

    # ============================================================
    # OCR HELPERS (Tesseract only — v3.0)
    # ============================================================

    def _ocr_pdf_page(self, pdf_path: str, page_num: int) -> str:
        """Single PDF page ko image me convert karke OCR lagao."""
        if not self.tesseract_available:
            return ""
        try:
            from pdf2image import convert_from_path
            images = convert_from_path(
                pdf_path,
                first_page=page_num + 1,
                last_page=page_num + 1,
                dpi=200,
            )
            if images:
                # 🛡️ FIX F-S5: Use TEMP_DIR (no hardcoded paths)
                img_path = str(TEMP_DIR / f"temp_page_{page_num}_{uuid.uuid4().hex[:6]}.png")
                images[0].save(img_path, 'PNG')
                try:
                    img = Image.open(img_path)
                    text = _tesseract_ocr_image(img, self.languages)
                finally:
                    _secure_cleanup(img_path)
                return text
        except ImportError:
            logger.warning("⚠️ pdf2image not installed. Skipping OCR for scanned pages.")
        except Exception as e:
            logger.error(f"❌ OCR page failed: {e}")
        return ""

    def _ocr_entire_pdf(self, pdf_path: str) -> str:
        """Pure OCR fallback - jab PyPDF2 bilkul kaam na kare."""
        if not self.tesseract_available:
            return ""
        try:
            from pdf2image import convert_from_path
            images = convert_from_path(pdf_path, dpi=200)
            all_text = ""
            for i, img in enumerate(images):
                # 🛡️ FIX F-S5: Use TEMP_DIR
                img_path = str(TEMP_DIR / f"temp_page_{i}_{uuid.uuid4().hex[:6]}.png")
                img.save(img_path, 'PNG')
                try:
                    page_text = _tesseract_ocr_image(img, self.languages)
                finally:
                    _secure_cleanup(img_path)
                if page_text:
                    all_text += f"\n--- Page {i+1} (OCR) ---\n{page_text}\n"
            return all_text
        except Exception as e:
            logger.error(f"❌ Full PDF OCR failed: {e}")
            return ""

    # ============================================================
    # TEXT CLEANING
    # ============================================================

    def _clean_text(self, text: str) -> str:
        """Extracted text ko clean karo."""
        if not text:
            return ""
        # Remove excessive whitespace but preserve paragraph breaks
        text = re.sub(r'[ \t]+', ' ', text)
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip()

    @staticmethod
    def get_supported_formats() -> dict:
        """Sab supported formats ki list return karo."""
        return SUPPORTED_FORMATS

    @staticmethod
    def get_accepted_extensions() -> list:
        """Streamlit file_uploader ke liye extensions list."""
        return ACCEPTED_EXTENSIONS
